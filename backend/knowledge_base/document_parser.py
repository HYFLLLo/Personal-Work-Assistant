"""
文档解析器
支持多种文档类型的解析和文本提取
"""

import os
import re
from typing import Dict, List, Optional, BinaryIO, Any
from pathlib import Path
import logging

from .models import DocumentType, DocumentStatus

logger = logging.getLogger(__name__)


class DocumentParser:
    """文档解析器类"""
    
    # 支持的文件扩展名映射
    SUPPORTED_EXTENSIONS = {
        '.txt': DocumentType.TXT,
        '.md': DocumentType.MD,
        '.doc': DocumentType.DOC,
        '.docx': DocumentType.DOCX,
        '.xls': DocumentType.XLS,
        '.xlsx': DocumentType.XLSX,
        '.ppt': DocumentType.PPT,
        '.pptx': DocumentType.PPTX,
        '.pdf': DocumentType.PDF,
    }
    
    # 最大文件大小 (50MB)
    MAX_FILE_SIZE = 50 * 1024 * 1024
    
    def __init__(self):
        self.parsers = {
            DocumentType.TXT: self._parse_text,
            DocumentType.MD: self._parse_text,
            DocumentType.DOCX: self._parse_docx,
            DocumentType.DOC: self._parse_doc,
            DocumentType.XLSX: self._parse_excel,
            DocumentType.XLS: self._parse_excel,
            DocumentType.PPTX: self._parse_pptx,
            DocumentType.PPT: self._parse_ppt,
            DocumentType.PDF: self._parse_pdf,
        }
    
    def get_document_type(self, filename: str) -> Optional[DocumentType]:
        """根据文件名获取文档类型"""
        ext = Path(filename).suffix.lower()
        return self.SUPPORTED_EXTENSIONS.get(ext)
    
    def is_supported(self, filename: str) -> bool:
        """检查文件类型是否支持"""
        return self.get_document_type(filename) is not None
    
    def parse(self, file_path: str) -> Dict[str, Any]:
        """
        解析文档
        
        Args:
            file_path: 文件路径
            
        Returns:
            Dict包含: content(文本内容), metadata(元数据), status(状态)
        """
        try:
            # 检查文件大小
            file_size = os.path.getsize(file_path)
            if file_size > self.MAX_FILE_SIZE:
                raise ValueError(f"文件大小超过限制 ({self.MAX_FILE_SIZE / 1024 / 1024}MB)")
            
            # 获取文档类型
            doc_type = self.get_document_type(file_path)
            if not doc_type:
                raise ValueError(f"不支持的文件类型: {Path(file_path).suffix}")
            
            # 调用对应的解析器
            parser = self.parsers.get(doc_type)
            if not parser:
                raise ValueError(f"未找到对应的解析器: {doc_type}")
            
            result = parser(file_path)
            result['status'] = DocumentStatus.COMPLETED
            result['file_size'] = file_size
            result['file_type'] = doc_type
            
            logger.info(f"文档解析成功: {file_path}, 类型: {doc_type}, 内容长度: {len(result.get('content', ''))}")
            
            return result
            
        except Exception as e:
            logger.error(f"文档解析失败: {file_path}, 错误: {str(e)}")
            return {
                'content': '',
                'metadata': {},
                'status': DocumentStatus.FAILED,
                'error_message': str(e),
                'file_size': os.path.getsize(file_path) if os.path.exists(file_path) else 0,
                'file_type': self.get_document_type(file_path)
            }
    
    def _parse_text(self, file_path: str) -> Dict[str, Any]:
        """解析文本文件"""
        encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']
        content = None
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                break
            except UnicodeDecodeError:
                continue
        
        if content is None:
            raise ValueError("无法识别文件编码")
        
        return {
            'content': content,
            'metadata': {
                'encoding': encoding,
                'line_count': len(content.split('\n'))
            }
        }
    
    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """解析Word文档 (.docx)"""
        try:
            from docx import Document
            doc = Document(file_path)
            
            # 提取段落文本
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            content = '\n'.join(paragraphs)
            
            # 提取表格内容
            tables_text = []
            for table in doc.tables:
                table_rows = []
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    table_rows.append(' | '.join(row_text))
                if table_rows:
                    tables_text.append('\n'.join(table_rows))
            
            if tables_text:
                content += '\n\n[表格内容]\n' + '\n\n'.join(tables_text)
            
            return {
                'content': content,
                'metadata': {
                    'paragraph_count': len(paragraphs),
                    'table_count': len(doc.tables),
                    'word_count': len(content)
                }
            }
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx")
        except Exception as e:
            raise ValueError(f"DOCX解析失败: {str(e)}")
    
    def _parse_doc(self, file_path: str) -> Dict[str, Any]:
        """解析旧版Word文档 (.doc)"""
        # .doc格式较复杂，建议使用antiword或转换为docx
        # 这里提供一个简化实现
        try:
            # 尝试使用textract
            import textract
            content = textract.process(file_path).decode('utf-8')
            return {
                'content': content,
                'metadata': {'word_count': len(content)}
            }
        except ImportError:
            logger.warning("textract未安装，尝试其他方法解析.doc文件")
            # 作为备选，尝试直接读取（可能会是乱码）
            try:
                with open(file_path, 'rb') as f:
                    raw_content = f.read()
                    # 尝试提取文本
                    content = raw_content.decode('utf-8', errors='ignore')
                    # 清理控制字符
                    content = re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f]', '', content)
                    return {
                        'content': content,
                        'metadata': {'note': '使用备用方法解析，可能不完整'}
                    }
            except Exception as e:
                raise ValueError(f"DOC解析失败: {str(e)}")
    
    def _parse_excel(self, file_path: str) -> Dict[str, Any]:
        """解析Excel文件"""
        try:
            import pandas as pd
            
            # 读取所有sheet
            xl_file = pd.ExcelFile(file_path)
            all_sheets_content = []
            
            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # 将DataFrame转为文本
                sheet_content = f"[Sheet: {sheet_name}]\n"
                sheet_content += df.to_string(index=False)
                all_sheets_content.append(sheet_content)
            
            content = '\n\n'.join(all_sheets_content)
            
            return {
                'content': content,
                'metadata': {
                    'sheet_count': len(xl_file.sheet_names),
                    'sheet_names': xl_file.sheet_names,
                    'row_count': sum(pd.read_excel(file_path, sheet_name=sheet).shape[0] 
                                   for sheet in xl_file.sheet_names)
                }
            }
        except ImportError:
            raise ImportError("请安装 pandas 和 openpyxl: pip install pandas openpyxl")
        except Exception as e:
            raise ValueError(f"Excel解析失败: {str(e)}")
    
    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """解析PowerPoint文件 (.pptx)"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            all_slides = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = [f"[Slide {i}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:
                    all_slides.append('\n'.join(slide_text))
            
            content = '\n\n'.join(all_slides)
            
            return {
                'content': content,
                'metadata': {
                    'slide_count': len(prs.slides),
                    'word_count': len(content)
                }
            }
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")
        except Exception as e:
            raise ValueError(f"PPTX解析失败: {str(e)}")
    
    def _parse_ppt(self, file_path: str) -> Dict[str, Any]:
        """解析旧版PowerPoint (.ppt)"""
        # .ppt格式较复杂，建议使用转换工具
        try:
            import textract
            content = textract.process(file_path).decode('utf-8')
            return {
                'content': content,
                'metadata': {'note': '使用textract解析'}
            }
        except ImportError:
            raise ValueError(".ppt文件需要textract库，请安装: pip install textract")
    
    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """解析PDF文件"""
        try:
            import PyPDF2
            
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                
                all_pages = []
                for i, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            all_pages.append(f"[Page {i + 1}]\n{page_text}")
                    except Exception as e:
                        logger.warning(f"PDF第{i+1}页解析失败: {str(e)}")
                
                content = '\n\n'.join(all_pages)
                
                return {
                    'content': content,
                    'metadata': {
                        'page_count': len(pdf_reader.pages),
                        'parsed_pages': len(all_pages),
                        'word_count': len(content)
                    }
                }
        except ImportError:
            raise ImportError("请安装 PyPDF2: pip install PyPDF2")
        except Exception as e:
            raise ValueError(f"PDF解析失败: {str(e)}")
    
    def chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 50) -> List[Dict[str, Any]]:
        """
        将文本分割成块
        
        Args:
            text: 原始文本
            chunk_size: 每块的最大字符数
            overlap: 块之间的重叠字符数
            
        Returns:
            文本块列表，每个块包含content, start_pos, end_pos
        """
        if not text:
            return []
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = min(start + chunk_size, text_length)
            
            # 如果不是最后一块，尝试在句子边界处分割
            if end < text_length:
                # 查找最近的句子结束符
                for sep in ['。', '！', '？', '.', '!', '?', '\n']:
                    pos = text.rfind(sep, start, end)
                    if pos != -1 and pos > start + chunk_size // 2:
                        end = pos + 1
                        break
            
            chunk_content = text[start:end].strip()
            if chunk_content:
                chunks.append({
                    'content': chunk_content,
                    'start_pos': start,
                    'end_pos': end
                })
            
            # 移动起始位置，考虑重叠
            start = end - overlap if end < text_length else end
        
        return chunks


# 全局解析器实例
document_parser = DocumentParser()
